from pathlib import Path
import math

import matplotlib.pyplot as plt
from matplotlib.patches import Circle, Rectangle
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parents[1]
ASSETS = ROOT / "generated_assets"
OUTPUT = ROOT / "三舵轮底盘控制答辩PPT.pptx"
RENDER_IMG = ROOT / "隙锋渲染图.png"
GAZEBO_IMG = ROOT / "gazebo仿真.png"


def ensure_dir(path: Path) -> None:
    path.mkdir(parents=True, exist_ok=True)


def make_chassis_diagram() -> Path:
    ensure_dir(ASSETS)
    out = ASSETS / "three_steer_layout.png"

    fig, ax = plt.subplots(figsize=(6, 6))
    ax.set_aspect("equal")

    body = Rectangle((-0.18, -0.18), 0.36, 0.36, fill=False, linewidth=2, edgecolor="#1f4e79")
    ax.add_patch(body)

    wheels = [
        ("front", 0.32, 0.0),
        ("left", -0.16, 0.277128),
        ("right", -0.16, -0.277128),
    ]
    for name, x, y in wheels:
        ax.add_patch(Circle((x, y), 0.07, fill=False, linewidth=2, edgecolor="#c55a11"))
        ax.plot([0, x], [0, y], linestyle="--", color="#7f7f7f", linewidth=1)
        ax.text(x, y + 0.1, name, ha="center", va="bottom", fontsize=11)
        ax.text(x, y - 0.12, f"({x:.2f}, {y:.2f})", ha="center", va="top", fontsize=9)

    ax.arrow(0, 0, 0.22, 0, head_width=0.025, head_length=0.03, fc="#2f5597", ec="#2f5597", linewidth=2)
    ax.arrow(0, 0, 0, 0.22, head_width=0.025, head_length=0.03, fc="#70ad47", ec="#70ad47", linewidth=2)
    ax.text(0.25, 0.0, "X", fontsize=12, color="#2f5597", va="center")
    ax.text(0.0, 0.25, "Y", fontsize=12, color="#70ad47", ha="center")
    ax.text(0, -0.28, "Three-steer chassis layout", ha="center", fontsize=13)

    ax.set_xlim(-0.5, 0.5)
    ax.set_ylim(-0.5, 0.5)
    ax.axis("off")
    plt.tight_layout()
    plt.savefig(out, dpi=220, bbox_inches="tight")
    plt.close(fig)
    return out


def set_run_font(run, size=20, bold=False, color=(0, 0, 0), name="Microsoft YaHei"):
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.name = name
    run.font.color.rgb = RGBColor(*color)


def add_title(slide, title, subtitle=None):
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(11.8), Inches(0.9))
    tf = title_box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run()
    r.text = title
    set_run_font(r, size=28, bold=True, color=(31, 78, 121))

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.65), Inches(1.15), Inches(11), Inches(0.45))
        tf = sub_box.text_frame
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        set_run_font(r, size=12, color=(96, 96, 96))


def add_bullets(slide, items, left=0.9, top=1.7, width=5.6, height=4.8, font_size=20):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.word_wrap = True
    tf.clear()
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.text = item
        p.level = 0
        p.space_after = Pt(10)
        for run in p.runs:
            set_run_font(run, size=font_size)
    return box


def add_section_box(slide, title, body_lines, left, top, width, height):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(242, 247, 252)
    shape.line.color.rgb = RGBColor(180, 198, 231)

    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_run_font(r, size=18, bold=True, color=(31, 78, 121))

    for line in body_lines:
        p = tf.add_paragraph()
        p.text = line
        p.level = 0
        p.space_before = Pt(4)
        p.space_after = Pt(2)
        for run in p.runs:
            set_run_font(run, size=14)


def add_picture_fit(slide, image_path, left, top, width, height):
    slide.shapes.add_picture(str(image_path), Inches(left), Inches(top), width=Inches(width), height=Inches(height))


def build_ppt():
    ensure_dir(ASSETS)
    layout_img = make_chassis_diagram()

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank = prs.slide_layouts[6]

    # Slide 1
    slide = prs.slides.add_slide(blank)
    bg = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(245, 248, 252)
    bg.line.fill.background()

    add_title(slide, "机器人移动底盘的全向运动控制方法研究与实现", "毕业设计答辩PPT")
    add_bullets(
        slide,
        [
            "学生：黄增祥",
            "学号：2223515366",
            "专业班级：自动化2202",
            "指导教师：张光华",
        ],
        left=0.9,
        top=2.0,
        width=4.4,
        height=2.2,
        font_size=22,
    )
    add_section_box(
        slide,
        "课题定位",
        [
            "以企业实际三舵轮底盘为对象",
            "主线是 STM32 底盘控制代码设计",
            "ROS 2 / Gazebo 作为辅助验证平台",
        ],
        left=0.9,
        top=4.7,
        width=4.9,
        height=1.7,
    )
    if RENDER_IMG.exists():
        add_picture_fit(slide, RENDER_IMG, 8.7, 0.9, 2.8, 6.1)
    else:
        add_picture_fit(slide, layout_img, 7.6, 4.1, 4.2, 2.8)

    # Slide 2
    slide = prs.slides.add_slide(blank)
    add_title(slide, "课题来源与任务要求")
    add_bullets(
        slide,
        [
            "选题来源：就业所在单位项目。",
            "题目名称：机器人移动底盘的全向运动控制方法研究与实现。",
            "课题背景聚焦三舵轮全向舵轮底盘的运动学模型解算、多轮协同控制与实时遥操作。",
            "培养目标是综合训练机器人运动学、实时控制系统设计和人机交互实现能力。",
        ],
    )
    add_section_box(
        slide,
        "任务书主要任务",
        [
            "1. 建立三舵轮底盘逆运动学模型",
            "2. 设计实时模型解算与协同控制方法",
            "3. 集成 USB 无线手柄实现实时遥操作",
            "4. 完成 MCU 架构、仿真与实物验证、论文撰写",
        ],
        left=7.0,
        top=1.9,
        width=5.5,
        height=3.1,
    )

    # Slide 3
    slide = prs.slides.add_slide(blank)
    add_title(slide, "课题背景与研究意义")
    add_bullets(
        slide,
        [
            "三舵轮底盘具有全向运动能力，适用于狭小空间和灵活机动场景。",
            "企业已提供实物平台与执行器控制接口，需要完成底盘层控制逻辑实现。",
            "在实物联调成本较高的情况下，先通过 ROS 2 + Gazebo 进行控制思路验证更稳妥。",
            "本课题重点不在电机通信协议，而在速度指令到舵角/轮速控制量的生成与实现。",
        ],
    )
    add_section_box(
        slide,
        "课题要求对齐",
        [
            "实时运动学解算模块",
            "USB 手柄驱动与控制映射模块",
            "多舵轮协同控制器",
            "仿真与实物平台效果演示",
        ],
        left=7.0,
        top=1.9,
        width=5.4,
        height=3.0,
    )

    # Slide 4
    slide = prs.slides.add_slide(blank)
    add_title(slide, "系统总体方案")
    add_section_box(slide, "输入层", ["手柄按键输入", "上层速度指令 cmd_vel"], 0.8, 1.7, 2.2, 1.4)
    add_section_box(slide, "控制层", ["X / Y / ω 目标速度生成", "速度平滑处理", "逆运动学解算"], 3.4, 1.7, 2.5, 1.8)
    add_section_box(slide, "执行层", ["三舵轮舵角命令", "三行走轮轮速命令", "执行器接口下发"], 6.4, 1.7, 2.5, 1.8)
    add_section_box(slide, "反馈层", ["电机角度/速度反馈", "里程计速度更新", "状态与故障检测"], 9.3, 1.7, 2.6, 1.8)
    add_bullets(
        slide,
        [
            "实物控制主线：STM32 程序负责把车体速度目标转换为 3 个舵轮角度和 3 个轮速命令。",
            "仿真验证主线：ROS 2 节点接收 cmd_vel，输出转向组和轮速组控制命令，验证运动关系是否合理。",
        ],
        left=0.9,
        top=4.2,
        width=11.5,
        height=2.0,
        font_size=18,
    )

    # Slide 5
    slide = prs.slides.add_slide(blank)
    add_title(slide, "企业平台与三舵轮结构")
    if RENDER_IMG.exists():
        add_picture_fit(slide, RENDER_IMG, 0.7, 1.5, 2.8, 5.4)
    add_picture_fit(slide, layout_img, 3.6, 1.7, 4.1, 4.6)
    add_bullets(
        slide,
        [
            "选题审核表中明确硬件条件已具备：企业现有完整实验平台及 USB 无线手柄。",
            "本课题在企业平台基础上，重点完成底盘层控制算法与程序实现。",
            "三舵轮控制对象可抽象为 3 个转向关节 + 3 个行走轮。",
        ],
        left=8.0,
        top=1.9,
        width=4.4,
        height=3.6,
        font_size=17,
    )

    # Slide 6
    slide = prs.slides.add_slide(blank)
    add_title(slide, "三舵轮底盘结构与运动学")
    add_picture_fit(slide, layout_img, 0.8, 1.7, 5.0, 4.6)
    add_bullets(
        slide,
        [
            "三舵轮几何位置来自仿真模型与控制代码：前轮 (0.32, 0.00)，左右后轮 (-0.16, ±0.277)。",
            "车体速度由 Vx、Vy、ω 描述，每个轮心速度由平动速度与角速度叠加得到。",
            "在控制实现中，需要进一步求取各舵轮目标角度与对应轮速。",
            "工程实现中还要考虑舵轮机械角度范围、补角解与轮速反向问题。",
        ],
        left=6.2,
        top=1.8,
        width=6.2,
        height=4.5,
        font_size=18,
    )

    # Slide 7
    slide = prs.slides.add_slide(blank)
    add_title(slide, "STM32 控制程序设计")
    add_section_box(
        slide,
        "关键代码模块",
        [
            "KinematicSteer3.c：三舵轮逆运动学与里程计",
            "JoyControl.c：手柄速度映射与平滑处理",
            "ActuatorKincoSteer3.c：执行器控制与状态机",
            "SafeMgmtTask.c：任务调度与安全相关处理",
        ],
        left=0.8,
        top=1.6,
        width=5.7,
        height=3.0,
    )
    add_section_box(
        slide,
        "已实现功能",
        [
            "舵角就近选取与补角反向",
            "舵角超限时角速度约束",
            "转向未到位时暂停行走",
            "实时角度/速度反馈与状态输出",
        ],
        left=6.8,
        top=1.6,
        width=5.6,
        height=3.0,
    )
    add_bullets(
        slide,
        [
            "实物控制不是直接设定底盘动作，而是先解算 3 个舵轮角度和 3 个轮速，再通过执行器接口下发。",
        ],
        left=0.9,
        top=5.2,
        width=11.3,
        height=0.9,
        font_size=18,
    )

    # Slide 8
    slide = prs.slides.add_slide(blank)
    add_title(slide, "遥控操作与速度平滑控制")
    add_bullets(
        slide,
        [
            "手柄输入可映射为底盘 X、Y、ω 三方向目标速度。",
            "程序中为 X / Y / ω 三轴分别建立了三角形加减速轨迹，减少速度突变带来的冲击。",
            "当目标速度反向切换时，先过零再追踪新目标，增强控制平顺性。",
            "速度平滑后的 now_x / now_y / now_w 最终进入底盘 WalkCtrl 接口，实现遥控到底盘运动的完整链路。",
        ],
        left=0.9,
        top=1.8,
        width=11.3,
        height=3.8,
        font_size=19,
    )
    add_section_box(
        slide,
        "控制意义",
        [
            "减小手柄输入抖动影响",
            "降低实物启动/变向冲击",
            "为三舵轮协调动作提供更平滑输入",
        ],
        left=7.7,
        top=4.8,
        width=4.5,
        height=1.5,
    )

    # Slide 9
    slide = prs.slides.add_slide(blank)
    add_title(slide, "ROS 2 / Gazebo 仿真平台")
    add_section_box(
        slide,
        "仿真组成",
        [
            "URDF/Xacro 建模：three_steer.urdf.xacro",
            "Gazebo 启动：three_steer_gazebo.launch.py",
            "控制器配置：controllers.yaml",
            "速度转换节点：cmd_vel_to_three_steer.py",
        ],
        left=0.8,
        top=1.6,
        width=5.8,
        height=2.9,
    )
    add_section_box(
        slide,
        "仿真作用",
        [
            "验证底盘几何参数是否合理",
            "验证 cmd_vel 到舵角/轮速的转换关系",
            "记录典型工况 bag 数据并绘图分析",
            "为实物测试提供先验参考",
        ],
        left=6.8,
        top=1.6,
        width=5.5,
        height=2.9,
    )
    add_bullets(
        slide,
        [
            "仿真在本课题中定位为“辅助验证手段”，主线仍然是底盘控制代码设计与实物联调。",
        ],
        left=0.9,
        top=5.0,
        width=6.0,
        height=0.9,
        font_size=18,
    )
    if GAZEBO_IMG.exists():
        add_picture_fit(slide, GAZEBO_IMG, 7.3, 4.6, 5.1, 2.2)

    # Slide 10
    slide = prs.slides.add_slide(blank)
    add_title(slide, "仿真结果：纯前进与纯横移")
    add_picture_fit(slide, ROOT / "bags" / "bag_forward" / "cmd_vel.png", 0.7, 1.5, 3.9, 2.2)
    add_picture_fit(slide, ROOT / "bags" / "bag_forward" / "steer_angles.png", 4.8, 1.5, 3.9, 2.2)
    add_picture_fit(slide, ROOT / "bags" / "bag_forward" / "wheel_speeds.png", 8.9, 1.5, 3.9, 2.2)
    add_picture_fit(slide, ROOT / "bags" / "bag_lateral" / "cmd_vel.png", 0.7, 4.1, 3.9, 2.2)
    add_picture_fit(slide, ROOT / "bags" / "bag_lateral" / "steer_angles.png", 4.8, 4.1, 3.9, 2.2)
    add_picture_fit(slide, ROOT / "bags" / "bag_lateral" / "wheel_speeds.png", 8.9, 4.1, 3.9, 2.2)

    # Slide 11
    slide = prs.slides.add_slide(blank)
    add_title(slide, "仿真结果：纯自转")
    add_picture_fit(slide, ROOT / "bags" / "bag_rotate" / "cmd_vel.png", 0.8, 1.6, 4.0, 2.6)
    add_picture_fit(slide, ROOT / "bags" / "bag_rotate" / "steer_angles.png", 4.7, 1.6, 4.0, 2.6)
    add_picture_fit(slide, ROOT / "bags" / "bag_rotate" / "wheel_speeds.png", 8.6, 1.6, 4.0, 2.6)
    add_bullets(
        slide,
        [
            "已完成前进、横移、自转三类典型工况下的 bag 记录与命令曲线绘制。",
            "结果能够直观反映输入速度与三舵轮控制命令之间的对应关系。",
        ],
        left=1.0,
        top=4.8,
        width=11.2,
        height=1.1,
        font_size=18,
    )

    # Slide 12
    slide = prs.slides.add_slide(blank)
    add_title(slide, "阶段性成果、存在问题与下一步计划")
    add_section_box(
        slide,
        "阶段性成果",
        [
            "完成三舵轮底盘结构分析与运动学梳理",
            "完成 STM32 侧底盘控制主链路阅读与实现",
            "完成手柄遥控与速度平滑控制接入",
            "完成 ROS 2 / Gazebo 仿真平台搭建与曲线绘制",
        ],
        left=0.8,
        top=1.6,
        width=3.9,
        height=3.6,
    )
    add_section_box(
        slide,
        "当前问题",
        [
            "舵角连续性与机械限位问题仍需进一步优化",
            "实物响应受摩擦、间隙和执行器延迟影响",
            "仿真与实物的一致性仍需更多测试支撑",
        ],
        left=4.9,
        top=1.6,
        width=3.7,
        height=3.6,
    )
    add_section_box(
        slide,
        "下一步计划",
        [
            "继续开展典型工况实物测试",
            "优化控制平滑性和舵轮协调性",
            "补充实验结果分析与论文撰写",
            "整理答辩材料与展示内容",
        ],
        left=8.8,
        top=1.6,
        width=3.7,
        height=3.6,
    )

    prs.save(OUTPUT)
    print(f"Saved PPT: {OUTPUT}")


if __name__ == "__main__":
    build_ppt()
